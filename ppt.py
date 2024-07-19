from dotenv import load_dotenv
load_dotenv()

import streamlit as st
import os
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches

genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))

model = genai.GenerativeModel('gemini-pro')

# Function for getting response from the Gemini Model
def gemini_res(question: str) -> str:
    response = model.generate_content(question)
    return response.text

# Function to create PowerPoint presentation
def create_ppt(content: str, file_name: str):
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Generated PowerPoint Presentation"
    subtitle.text = "By Vaibhav"
    
    # Content Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "User Input and Response"
    body.text = content
    
    # Save the presentation
    prs.save(file_name)

# Streamlit App configuration
st.set_page_config(page_title='Write topic name for Generating PPT')

st.header('Topic name for Generating PPT')

# Input field
user_input = st.text_input('Input:', key='input')

# Submit button
submit = st.button('Send your Question and get Answer')

# When submit button is clicked
if submit:
    responses = gemini_res(user_input)
    st.subheader('Your Answer is')
    st.write(responses)
    
    # Combine user input and response for the PowerPoint content
    ppt_content = f"User Input: {user_input}\n\nResponse: {responses}"
    
    # Create and save the PowerPoint presentation
    ppt_file = "output_presentation.pptx"
    create_ppt(ppt_content, ppt_file)
    
    st.success(f"PowerPoint presentation saved as {ppt_file}")